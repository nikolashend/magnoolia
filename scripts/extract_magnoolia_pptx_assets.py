"""
Phase 28 — PPTX Asset Extractor
Extracts images and text from Magnoolia PPTX files.
Usage: python scripts/extract_magnoolia_pptx_assets.py
"""
import os
import sys
import json
import shutil
import zipfile
from pathlib import Path

ROOT = Path(__file__).parent.parent
IMPORT_DIR = ROOT / "materials" / "onedrive-import" / "phase28"
OUT_BASE = ROOT / "storage" / "app" / "magnoolia" / "phase28" / "pptx-extracted"
PUBLIC_OUT = ROOT / "public" / "assets" / "magnoolia" / "sisedisain" / "pptx"
REPORT_PATH = ROOT / "docs" / "magnoolia-phase28-pptx-extraction-report.md"

OUT_BASE.mkdir(parents=True, exist_ok=True)
PUBLIC_OUT.mkdir(parents=True, exist_ok=True)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("WARNING: python-pptx not available, falling back to ZIP extraction only")

def extract_text_from_shape(shape):
    """Extract all text from a shape recursively."""
    texts = []
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
        if shape.shape_type == 6:  # GROUP
            for s in shape.shapes:
                texts.extend(extract_text_from_shape(s))
    except Exception:
        pass
    return texts

def process_pptx(pptx_path, out_dir, public_dir):
    pptx_name = pptx_path.stem
    slide_inventory = []
    media_count = 0

    print(f"\n{'='*60}")
    print(f"Processing: {pptx_path.name}")
    print(f"{'='*60}")

    # Step 1: Extract all media via ZIP
    media_dir = out_dir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    extracted_media = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            for name in z.namelist():
                if name.startswith('ppt/media/') and name != 'ppt/media/':
                    ext = Path(name).suffix.lower()
                    if ext in ['.jpg', '.jpeg', '.png', '.webp', '.gif', '.bmp', '.tiff']:
                        fname = Path(name).name
                        dest = media_dir / fname
                        with z.open(name) as src, open(dest, 'wb') as dst:
                            dst.write(src.read())
                        size_kb = dest.stat().st_size // 1024
                        extracted_media.append({'file': fname, 'path': str(dest), 'size_kb': size_kb})
                        media_count += 1
        print(f"  Extracted {media_count} media files from ZIP")
    except Exception as e:
        print(f"  ZIP extraction error: {e}")

    # Step 2: Parse slides with python-pptx if available
    if HAS_PPTX:
        try:
            prs = Presentation(str(pptx_path))
            print(f"  Slides: {len(prs.slides)}")

            for i, slide in enumerate(prs.slides, 1):
                slide_info = {
                    'slide': i,
                    'title': '',
                    'texts': [],
                    'image_count': 0,
                }
                all_texts = []

                for shape in slide.shapes:
                    # Get title
                    if shape.shape_type == 13:  # PICTURE
                        slide_info['image_count'] += 1
                    elif hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                        if shape.placeholder_format.idx == 0:  # Title
                            try:
                                slide_info['title'] = shape.text_frame.text.strip()
                            except Exception:
                                pass
                    texts = extract_text_from_shape(shape)
                    all_texts.extend(texts)

                slide_info['texts'] = [t for t in all_texts if len(t) > 2]
                slide_inventory.append(slide_info)

                if slide_info['title'] or slide_info['texts']:
                    print(f"  Slide {i:2d}: {slide_info['title'][:60] or '(no title)'} | {len(slide_info['texts'])} text blocks | {slide_info['image_count']} images")

        except Exception as e:
            print(f"  python-pptx parse error: {e}")

    # Step 3: Copy usable images to public output
    copied_public = []
    min_size_kb = 20  # Skip tiny icons/decorations
    for m in extracted_media:
        if m['size_kb'] >= min_size_kb:
            src = Path(m['path'])
            dst = public_dir / src.name
            shutil.copy2(src, dst)
            copied_public.append(m['file'])

    print(f"  Copied {len(copied_public)} usable images (>= {min_size_kb}KB) to public")

    return {
        'pptx': pptx_path.name,
        'slides': len(slide_inventory),
        'media_total': media_count,
        'media_public': len(copied_public),
        'slide_inventory': slide_inventory,
        'media': extracted_media,
        'public_files': copied_public,
    }

def main():
    results = []

    pptx_files = [
        IMPORT_DIR / "Magnoolia  kodud Prestige Sisedisain.pptx",
        IMPORT_DIR / "Magnoolia kodud Prestige.pptx",
    ]

    for pptx_path in pptx_files:
        if not pptx_path.exists():
            print(f"NOT FOUND: {pptx_path}")
            continue

        out_dir = OUT_BASE / pptx_path.stem.replace(' ', '_').replace('  ', '_')
        out_dir.mkdir(parents=True, exist_ok=True)
        public_dir = PUBLIC_OUT / pptx_path.stem.replace(' ', '_').replace('  ', '_')
        public_dir.mkdir(parents=True, exist_ok=True)

        result = process_pptx(pptx_path, out_dir, public_dir)
        results.append(result)

    # Write inventory JSON
    inv_path = OUT_BASE / "inventory.json"
    with open(inv_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    print(f"\nInventory written: {inv_path}")

    # Write report
    lines = ["# Phase 28 PPTX Extraction Report", "", f"Generated: {__import__('datetime').datetime.now().isoformat()}", ""]
    for r in results:
        lines += [
            f"## {r['pptx']}",
            "",
            f"- Slides: {r['slides']}",
            f"- Media files extracted: {r['media_total']}",
            f"- Usable images copied to public: {r['media_public']}",
            "",
            "### Slide inventory",
            "",
            "| Slide | Title | Texts | Images |",
            "|-------|-------|-------|--------|",
        ]
        for s in r.get('slide_inventory', []):
            title = (s['title'] or '—')[:60]
            texts = len(s['texts'])
            lines.append(f"| {s['slide']} | {title} | {texts} | {s['image_count']} |")
        lines += ["", "### Public images copied", ""]
        for f in r.get('public_files', []):
            m = next((x for x in r['media'] if x['file'] == f), {})
            lines.append(f"- `{f}` ({m.get('size_kb', '?')} KB)")
        lines.append("")

    REPORT_PATH.write_text('\n'.join(lines), encoding='utf-8')
    print(f"Report: {REPORT_PATH}")

    # Summary
    total_media = sum(r['media_total'] for r in results)
    total_public = sum(r['media_public'] for r in results)
    print(f"\n{'='*60}")
    print(f"TOTAL: {len(results)} PPTX processed, {total_media} media files, {total_public} copied to public")
    print(f"{'='*60}\n")

if __name__ == '__main__':
    main()
